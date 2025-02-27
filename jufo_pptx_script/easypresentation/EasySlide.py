from pptx.dml.color import RGBColor

from jufo_pptx_script.easypresentation.EasyTextformatter import EasyTextformatter, EasyTextformatterList
from jufo_pptx_script.templater.Template import Template
from jufo_pptx_script.templater.TemplateApplier import ImageParseResult
from pptx.slide import Slide
import io
from jufo_pptx_script.easypresentation.EasyPresentation import EasyPresentation
from PIL import Image
from pptx.shapes.placeholder import SlidePlaceholder, PicturePlaceholder
from pptx.shapes.autoshape import Shape as AutoShape
from typing import Any, Literal
import sys


def _clamp(value, min_value, max_value):
    return max(min_value, min(value, max_value))


class EasySlide:

    # TODO: Maybe implement a cache for the __get_placeholders

    def __init__(self, slide: Slide, pptx: EasyPresentation):
        self.__slide = slide
        self.__pptx = pptx

    def __get_placeholders(self):
        """
        Create a list of 2-item tuples.
        First item being the new instance of the placeholders in the current slide
        Second item being the link to the original placeholder
        """
        ls = []

        # Iterates over the master-layout
        for io, pl_master in enumerate(self.__slide.slide_layout.placeholders):
            # Iterates over all new instances of placeholders
            for ii, pl_new in enumerate(self.__slide.placeholders):
                # Checks that they are equal
                if pl_master.placeholder_format.idx == pl_new.placeholder_format.idx:
                    ls.append((pl_new, pl_master))

        return ls

    def get_placeholder_types(self) -> dict[str, Literal["Picture"] or Literal["Text"] or None]:
        types = dict()

        placeholders = self.__get_placeholders()

        for i, (pln, plm) in enumerate(self.__get_placeholders()):
            # Gets the placeholder text on the original placeholder
            text = plm.text

            if isinstance(pln, SlidePlaceholder):
                types[text] = "Text"
            elif isinstance(pln, PicturePlaceholder):
                types[text] = "Picture"

        return types

    def apply_template(self, template: Template, data: dict[str, Any]):

        for shape in self.__slide.shapes:
            if isinstance(shape, AutoShape):
                self.__update_text_placeholder(data, shape.text, shape, template)


        for i, (pln, plm) in enumerate(self.__get_placeholders()):
            # Gets the placeholder text (Which is not accessible otherwise)
            text = plm.text

            if isinstance(pln, SlidePlaceholder):
                self.__update_text_placeholder(data, text, pln, template)
            if isinstance(pln, PicturePlaceholder):
                self.__update_image_placeholder(data, text, pln, template)

    # region Internal template apply mechanism

    @property
    def _image_cache(self):
        return self.__pptx._image_cache

    def _load_image(self, request: ImageParseResult):
        cache = self.__pptx._image_cache

        img_bytes = cache.get(request.raw) if cache is not None else None

        if img_bytes is None:
            img = Image.open(request.result)

            # Scales the image if required
            if request.scale is not None:
                # Checks for a percentage scale
                if type(request.scale) is int:
                    request.scale = _clamp(request.scale, 1, 100)
                    img = img.resize((int(img.width * request.scale / 100), int(img.height * request.scale / 100)))
                # It's a pixel scale
                else:
                    request.scale = (
                        _clamp(request.scale[0], 1, 99999),
                        _clamp(request.scale[1], 1, 99999)
                    )

                    img = img.resize(request.scale)

            img_bytes = io.BytesIO()
            img.save(img_bytes, img.format if img.format is not None else "jpeg")

            if cache is not None:
                cache[request.raw] = img_bytes

        img_bytes.seek(0)
        return img_bytes

    def __update_text_placeholder(self, data: dict[str, Any], raw_text: str,
                                  element: SlidePlaceholder or AutoShape, template: Template):
        res = self.__pptx._template.parse(template, data, raw_text)

        # Gets the text-frame (Used to apply formatting and texts)
        frame = element.text_frame

        for textList in res:

            elements = textList.elements

            for part in elements:

                # Ensures any none-compliant data types are converted to a string
                if not isinstance(part, EasyTextformatter):
                    part = str(part)

                # If only a text is given, apply it
                if type(part) is str:
                    frame.paragraphs[0].add_run().text = part
                    continue

                # Applies the formatted text
                part._apply_to_frame(frame)


    def __update_image_placeholder(self, data: dict[str, Any], raw_text: str,
                                   placeholder: PicturePlaceholder, template: Template):
        # Resolves any templates inside the image infos
        res = self.__pptx._template.parse_with_image_properties(template, data, raw_text)

        # Ensures that specific paths dont need to be loaded
        if res.result == 'IGNORE_FLAG':
            return

        try:
            img_as_bytes = self._load_image(res)
            placeholder.insert_picture(img_as_bytes)
        except FileNotFoundError:
            print(f"\rImage at '{res.result}' was not found, skipping...",file=sys.stderr)
        except Exception as err:
            raise ValueError(f"Image at '{res.result}' couldn't be loaded: {err}")


    # endregion
