from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from jufo_pptx_script.templater.TemplateApplier import TemplateApplier
import io
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml

class EasyPresentation:

    def __init__(self, path: str, use_image_cache: bool = False):
        self.__presentation: PresentationType = Presentation(path)
        self._template = TemplateApplier()
        self._image_cache: {str: io.BytesIO} or None = {} if use_image_cache else None

    def new_slide_from_layout(self, master: str, layout: str):
        from jufo_pptx_script.easypresentation.EasySlide import EasySlide

        slide_layout = self.__get_slidelayout_by_name(master, layout)

        new_slide = self.__presentation.slides.add_slide(slide_layout)
        return EasySlide(new_slide, self)

    def find_slide_by_note(self, note: str):
        from jufo_pptx_script.easypresentation.EasySlide import EasySlide

        for slide in self.__presentation.slides:
            if slide.notes_slide.notes_text_frame.text == note:
                return EasySlide(slide, self)

        raise FileNotFoundError(f"Slide with note '{note}' could not be found.")

    def save_to(self, path: str):
        self.__presentation.save(path)

    def __get_slidelayout_by_name(self, master: str, layout: str):
        found_master = None

        for slidemaster in self.__presentation.slide_masters:
            theme = slidemaster.part.part_related_by(RT.THEME)
            slide_name = parse_xml(theme.blob).get("name")
            if slide_name == master:
                found_master = slidemaster
                break

        if found_master is None:
            raise ValueError(f"The Slidemaster '{master}' couldn't be found but was requested.")

        for lay in found_master.slide_layouts:
            if lay.name == layout:
                return lay

        raise ValueError(f"The Layout '{layout}' couldn't be found on Slidemaster '{master}' but was requested.")